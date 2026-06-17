from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn as QN
import subprocess, re

TF='#C9E7B8';TB='#7FA968';AF='#D7ECF7';AB='#4E97BE';BF='#FDF3C4';BB='#C9A227'
# ---------- tech model data (shortened labels) ----------
LAB={'acm':'Access Control\\nManagement','jams':'Journey & Account\\nMgmt System','fms':'Fare Management\\nSystem','web':'ArchiTrain\\nWebsite','bank':'Bank Payment\\nService (ext)','do_tap':'Tap Event\\nData','do_acct':'Go Card\\nAccount Data','do_fz':'Fare & Zone\\nData','do_tic':'Digital Ticket','n_si':'Station\\nInfrastructure','n_be':'Journey & Account\\nMgmt Backend','n_as':'Application\\nServer','d_gate':'Station Gate','d_nfc':'NFC Reader\\n(embedded)','d_pnfc':'Platform\\nNFC Reader','d_card':'Go Card','d_wallet':'Mobile Wallet /\\nBank Card','ss_db':'Database (DBMS)','ss_fw':'Gate Firmware /\\nAccess Ctrl Client','if_nfc':'NFC Interface','sv_store':'Data Storage\\nService','sv_nfc':'NFC Proximity\\nComm Service','fn_store':'Store Journey\\n& Account Data','fn_hs':'NFC Handshake\\nProcessing','ar_db':'Journey & Account\\nDatabase','ar_web':'Website\\nDeployment','net':'Private Network','p_nfc':'NFC Handshake\\n(Path)','p_link':'Gate-Backend\\nLink (Path)'}
APP={'acm','jams','fms','web','bank','do_tap','do_acct','do_fz','do_tic'}
E=[('comp','n_si','d_gate'),('comp','n_si','d_pnfc'),('comp','d_gate','d_nfc'),('comp','d_gate','ss_fw'),('comp','n_be','ss_db'),('comp','n_be','ar_db'),('comp','n_as','ar_web'),('comp','d_nfc','if_nfc'),('assign','d_gate','acm'),('assign','n_be','jams'),('assign','n_be','fms'),('assign','n_as','web'),('assign','ss_db','fn_store'),('assign','ss_fw','fn_hs'),('real','n_be','sv_store'),('real','d_nfc','sv_nfc'),('real','d_pnfc','sv_nfc'),('real','ar_web','web'),('real','ar_db','do_tap'),('real','ar_db','do_acct'),('real','ar_db','do_fz'),('real','ar_db','do_tic'),('real','net','p_link'),('serving','sv_store','jams'),('serving','sv_store','fms'),('serving','sv_nfc','acm'),('serving','bank','fms'),('flow','fn_hs','fn_store'),('access','fn_store','ar_db'),('assoc','net','n_si'),('assoc','net','n_be'),('assoc','net','n_as'),('assoc','p_nfc','d_nfc'),('assoc','p_nfc','d_pnfc'),('assoc','p_nfc','d_card'),('assoc','p_nfc','d_wallet'),('assoc','p_link','d_gate'),('assoc','p_link','n_be')]
def tnode(k):
    fc=AF if k in APP else TF; ec=AB if k in APP else TB
    return f'"{k}"[label="{LAB[k]}",shape=box,style="rounded,filled",fillcolor="{fc}",color="{ec}",fontname="Helvetica-Bold",fontsize=13,width=2.3,height=0.95,fixedsize=true];'
d=['digraph G{rankdir=BT;ranksep=0.8;nodesep=0.6;','edge[fontname="Helvetica",fontsize=11];']
for k in LAB: d.append(tnode(k))
for t,s,g in E: d.append(f'"{s}"->"{g}"[label="{t}"];')
d.append('}'); open('/tmp/poster_tech_full.dot','w').write('\n'.join(d))

# ---------- abstract model data ----------
ABS={'b-passenger':('Passenger','b'),'b-css':('Customer\\nSupport Staff','b'),'b-controller':('Operations\\nController','b'),'b-trainop':('Train\\nOperator','b'),'b-malf':('Malfunction\\nManagement','b'),
'b-svc-ticket':('Ticketing &\\nFare Service','b'),'b-svc-track':('Train Operation &\\nTracking Service','b'),'b-svc-info':('Passenger\\nInformation Service','b'),
'b-proc-journey':('Handle Journeys\\n& Fares','b'),'b-proc-local':('Localize &\\nMonitor Trains','b'),'b-proc-delay':('Handle Delays\\n& Disruptions','b'),
'b-ev-tapoff':('Tap-off Event\\nCreated','b'),'b-ev-threshold':('Delay exceeds\\nthreshold','b'),
'a-website':('ArchiTrain\\nWebsite','a'),'a-jams':('Journey & Account\\nMgmt System','a'),'a-fms':('Fare Management\\nSystem','a'),'a-acm':('Access Control\\nManagement','a'),'a-occ':('Operational\\nControl Center','a'),'a-pis':('Passenger\\nInformation System','a'),'a-tms':('Train Management\\nSystem','a'),'a-dhs':('Disruption\\nHandling System','a'),'a-lccu-sw':('LCCU Control\\nSoftware','a'),
'a-svc-fare':('Fare Calculation\\nService','a'),'a-svc-track':('Train Tracking\\nService','a'),'a-svc-delay':('Delay Notification\\nService','a'),'a-svc-gate':('Gate Control\\nService','a'),'a-svc-bank':('Bank Payment\\nService (ext)','a'),
't-station':('Station\\nInfrastructure','t'),'t-jams-node':('Journey & Account\\nMgmt Backend','t'),'t-appserver':('Application\\nServer','t'),'t-lccu':('LCCU\\n(On-board)','t'),'t-backend':('Central Railway\\nBackend','t'),'t-net-priv':('Private Network','t'),'t-net-radio':('Railway Radio','t')}
ABS_E=[('serv','b-svc-ticket','b-passenger'),('serv','b-svc-track','b-passenger'),('serv','b-svc-info','b-passenger'),
('real','b-proc-journey','b-svc-ticket'),('real','b-proc-local','b-svc-track'),('real','b-proc-delay','b-svc-info'),
('assign','b-passenger','b-proc-journey'),('assign','b-css','b-proc-journey'),('assign','b-controller','b-proc-delay'),
('aggr','b-malf','b-controller'),('aggr','b-malf','b-trainop'),('assign','b-malf','b-proc-delay'),
('trig','b-ev-tapoff','b-proc-journey'),('trig','b-ev-threshold','b-proc-delay'),
('real','a-fms','a-svc-fare'),('serv','a-svc-fare','b-proc-journey'),('real','a-acm','a-svc-gate'),('serv','a-svc-gate','b-proc-journey'),
('real','a-occ','a-svc-track'),('serv','a-svc-track','b-proc-local'),('real','a-occ','a-svc-delay'),('serv','a-svc-delay','b-proc-delay'),
('aggr','a-dhs','a-occ'),('aggr','a-dhs','a-pis'),('aggr','a-dhs','a-tms'),
('flow','a-lccu-sw','a-occ'),('flow','a-occ','a-pis'),('flow','a-occ','a-tms'),('flow','a-pis','a-website'),
('serv','a-jams','a-fms'),('serv','a-jams','a-website'),('serv','a-svc-bank','a-fms'),
('assign','t-station','a-acm'),('assign','t-backend','a-occ'),('assign','t-backend','a-pis'),('assign','t-backend','a-tms'),('assign','t-lccu','a-lccu-sw'),('assign','t-appserver','a-website'),('assign','t-jams-node','a-jams'),('assign','t-jams-node','a-fms'),
('assoc','t-net-radio','t-lccu'),('assoc','t-net-radio','t-backend'),('assoc','t-net-priv','t-station'),('assoc','t-net-priv','t-backend'),('assoc','t-net-priv','t-jams-node')]
def anode(k):
    nm,ly=ABS[k]; fc={'b':BF,'a':AF,'t':TF}[ly]; ec={'b':BB,'a':AB,'t':TB}[ly]
    return f'"{k}"[label="{nm}",shape=box,style="rounded,filled",fillcolor="{fc}",color="{ec}",fontname="Helvetica-Bold",fontsize=13,width=2.3,height=0.95,fixedsize=true];'
d=['digraph G{rankdir=BT;ranksep=0.9;nodesep=0.55;','edge[fontname="Helvetica",fontsize=11];']
for k in ABS: d.append(anode(k))
for t,s,g in ABS_E: d.append(f'"{s}"->"{g}"[label="{t}"];')
d.append('}'); open('/tmp/poster_abstract.dot','w').write('\n'.join(d))

# ---------- poster pptx ----------
ARROW={'comp':(None,'diamond'),'assign':('triangle',None),'real':('triangle',None),'serving':('stealth',None),'serv':('stealth',None),'flow':('stealth',None),'access':('arrow',None),'trig':('triangle',None),'aggr':(None,'diamond'),'assoc':(None,None)}
DASHED={'assoc':'dash','serving':'sysDot','serv':'sysDot','access':'sysDot','flow':'dash'}
prs=Presentation(); prs.slide_width=Inches(48); prs.slide_height=Inches(27); blank=prs.slide_layouts[6]
def rect(s,l,t,w,h,fill,line=None,lw=1.0):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid();sp.fill.fore_color.rgb=fill
    sp.line.color.rgb=line; sp.line.width=Pt(lw); sp.shadow.inherit=False
    try: sp.adjustments[0]=0.08
    except: pass
    return sp
def para(tf,txt,size,color,bold=False,center=True):
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    r=p.add_run(); r.text=txt; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name='Arial'
def set_ends(c,head,tail):
    ln=c.line._get_or_add_ln()
    if head: ln.append(ln.makeelement(QN('a:headEnd'),{'type':head,'w':'med','len':'med'}))
    if tail: ln.append(ln.makeelement(QN('a:tailEnd'),{'type':tail,'w':'med','len':'med'}))
def set_dash(c,v):
    ln=c.line._get_or_add_ln(); ln.append(ln.makeelement(QN('a:prstDash'),{'val':v}))
def poster(title,dotpath):
    s=prs.slides.add_slide(blank)
    txt=subprocess.run(['dot','-Tplain',dotpath],capture_output=True,text=True).stdout
    nodes={};edges=[];GW=GH=1.0
    for line in txt.splitlines():
        if line.startswith('graph '):
            p=line.split();GW=float(p[2]);GH=float(p[3])
        elif line.startswith('node '):
            m=re.match(r'node (\S+) (\S+) (\S+) (\S+) (\S+) "(.*?)" \S+ \S+ (\S+) (\S+)',line)
            if m: nodes[m.group(1)]=dict(cx=float(m.group(2)),cy=float(m.group(3)),w=float(m.group(4)),h=float(m.group(5)),lab=m.group(6).replace('\\n','\n'),bd=m.group(7),fl=m.group(8))
        elif line.startswith('edge '):
            p=line.split();n=int(p[3]);rest=p[4+2*n:]
            rt=rest[0] if rest else '';lx=float(rest[1]) if len(rest)>2 else 0;ly=float(rest[2]) if len(rest)>2 else 0
            edges.append((p[1],p[2],rt,lx,ly))
    M=1.2; bw=48-2*M; bh=27-2*M
    sc=min(bw/GW,bh/GH); offx=M+(bw-GW*sc)/2; offy=M+(bh-GH*sc)/2
    nf=13*sc; ef=10*sc
    def place(nd): return offx+(nd['cx']-nd['w']/2)*sc, offy+(GH-nd['cy']-nd['h']/2)*sc, nd['w']*sc, nd['h']*sc
    shp={}
    for nid,nd in nodes.items():
        x,y,w,h=place(nd)
        r=rect(s,x,y,w,h,RGBColor.from_string(nd['fl'][1:]),RGBColor.from_string(nd['bd'][1:]),lw=1.25)
        tf=r.text_frame;tf.word_wrap=True;tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        tf.margin_left=Pt(2);tf.margin_right=Pt(2);tf.margin_top=Pt(0);tf.margin_bottom=Pt(0)
        p=tf.paragraphs[0];p.alignment=PP_ALIGN.CENTER;p.line_spacing=0.95
        rr=p.add_run();rr.text=nd['lab'];rr.font.size=Pt(nf);rr.font.bold=False;rr.font.color.rgb=RGBColor(0x22,0x26,0x2B);rr.font.name='Arial'
        shp[nid]=r
    for tail,head,rt,lx,ly in edges:
        if tail in shp and head in shp:
            a=shp[tail];b=shp[head]
            c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,a.left+a.width//2,a.top+a.height//2,b.left+b.width//2,b.top+b.height//2)
            c.line.color.rgb=RGBColor(0x70,0x7A,0x85);c.line.width=Pt(1.25);c.shadow.inherit=False
            he,te=ARROW.get(rt,('arrow',None)); set_ends(c,he,te)
            if rt in DASHED: set_dash(c,DASHED[rt])
            try: c.begin_connect(a,0);c.end_connect(b,0)
            except Exception: pass
            if rt:
                ex=offx+lx*sc; ey=offy+(GH-ly)*sc; lw_=1.4; lh_=0.5
                lb=s.shapes.add_textbox(Inches(ex-lw_/2),Inches(ey-lh_/2),Inches(lw_),Inches(lh_))
                lb.fill.solid();lb.fill.fore_color.rgb=RGBColor(0xFF,0xFF,0xFF);lb.line.fill.background();lb.shadow.inherit=False
                t2=lb.text_frame;t2.word_wrap=False;t2.margin_left=Pt(0);t2.margin_right=Pt(0);t2.margin_top=Pt(0);t2.margin_bottom=Pt(0)
                pp=t2.paragraphs[0];pp.alignment=PP_ALIGN.CENTER
                rr=pp.add_run();rr.text=rt;rr.font.size=Pt(ef*0.8);rr.font.color.rgb=RGBColor(0x6A,0x72,0x7C);rr.font.name='Arial'
    # corner label
    tbx=s.shapes.add_textbox(Inches(0.6),Inches(0.4),Inches(20),Inches(1.0))
    para(tbx.text_frame,title,30,RGBColor(0xB0,0xB8,0xC2),bold=True,center=False)

posters=[('Technology layer — full model','/tmp/poster_tech_full.dot'),
('Technology — Perspective 1: Tap & gate','/tmp/tech_v1_tap.dot'),
('Technology — Perspective 2: Communication (P01)','/tmp/tech_v2_comm.dot'),
('Technology — Perspective 3: Backend storage & data','/tmp/tech_v3_storage.dot'),
('Technology — Perspective 4: Self-service & settlement','/tmp/tech_v4_self.dot'),
('Abstract model — whole enterprise','/tmp/poster_abstract.dot')]
for t,dp in posters: poster(t,dp)
prs.save('/Users/ziway/Downloads/ArchiTrain_Diagrams_Posters.pptx')
print("posters saved:",len(prs.slides._sldIdLst))
