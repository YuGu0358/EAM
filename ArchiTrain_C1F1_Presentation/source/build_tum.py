from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn as QN
import subprocess

TUMB=RGBColor(0x00,0x65,0xBD); TUMD=RGBColor(0x07,0x2A,0x4A)
BLACK=RGBColor(0x00,0x00,0x00); GREY=RGBColor(0x59,0x59,0x59); MUTE=RGBColor(0x8A,0x93,0x9E)
DARK=RGBColor(0x26,0x2A,0x30); WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xC8,0xD2,0xDE); ORANGE=RGBColor(0xC5,0x5A,0x11)
BIZ=RGBColor(0xFD,0xF3,0xC4);BIZ_B=RGBColor(0xC9,0xA2,0x27)
APPC=RGBColor(0xDB,0xEE,0xF8);APP_B=RGBColor(0x4E,0x97,0xBE)
TEC=RGBColor(0xDC,0xF0,0xCD);TEC_B=RGBColor(0x7F,0xA9,0x68)
DASH=RGBColor(0xB8,0xC2,0xCE)
F="Arial"
FOOT="C1F1G_ | ArchiTrain | Ticketing & Train Journey Management"

from pptx.oxml.ns import qn as _qn
prs=Presentation('/tmp/TUM_template.pptx')
SW=prs.slide_width/914400; SH=prs.slide_height/914400
# delete example slides (drop relationships so old parts are not re-written)
ids=prs.slides._sldIdLst
for sid in list(ids):
    rId=sid.get(_qn('r:id'))
    try: prs.part.drop_rel(rId)
    except Exception: pass
    ids.remove(sid)
# layout map
L={}
for m in prs.slide_masters:
    for lay in m.slide_layouts: L.setdefault(lay.name, lay)

def dims(p):
    o=subprocess.run(['sips','-g','pixelWidth','-g','pixelHeight',p],capture_output=True,text=True).stdout
    w=int([l for l in o.splitlines() if 'pixelWidth' in l][0].split(':')[1]);h=int([l for l in o.splitlines() if 'pixelHeight' in l][0].split(':')[1]);return w/h
def tb(s,l,t,w,h):
    b=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h));b.text_frame.word_wrap=True;return b
def para(tf,segs,size,space=3,first=False,align=PP_ALIGN.LEFT,ls=None):
    p=tf.paragraphs[0] if first else tf.add_paragraph();p.alignment=align;p.space_after=Pt(space)
    if ls:p.line_spacing=ls
    for txt,bold,color in segs:
        r=p.add_run();r.text=txt;r.font.size=Pt(size);r.font.bold=bold;r.font.color.rgb=color;r.font.name=F
    return p
def rect(s,l,t,w,h,fill,line=None,rounded=False,lw=1.0,dash=False):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid();shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else:
        shp.line.color.rgb=line;shp.line.width=Pt(lw)
        if dash:
            from pptx.oxml.ns import qn
            ln=shp.line._get_or_add_ln(); d=ln.makeelement(qn('a:prstDash'),{'val':'dash'}); ln.append(d)
    shp.shadow.inherit=False
    if rounded:
        try: shp.adjustments[0]=0.07
        except: pass
    return shp
def pic(s,p,top,height):
    ar=dims(p);w=height*ar;s.shapes.add_picture(p,Inches((SW-w)/2),Inches(top),height=Inches(height));return w
def pic_fit(s,p,x,y,maxw,maxh):
    ar=dims(p); w=maxw; h=w/ar
    if h>maxh: h=maxh; w=h*ar
    s.shapes.add_picture(p,Inches(x+(maxw-w)/2),Inches(y),width=Inches(w)); return w,h
CHIP=RGBColor(0xEC,0xF1,0xF6)
def stepper(s,active):
    steps=["Overview","1 Tap & Gate","2 Communication","3 Storage","4 Self-service"]
    x=0.55; cw=1.74; gap=0.05; h=0.3; y=1.62
    for i,st in enumerate(steps):
        on=(i==active); rect(s,x,y,cw,h,TUMB if on else CHIP,rounded=True)
        para(tb(s,x,y+0.045,cw,0.24).text_frame,[(st,True,WHITE if on else GREY)],8.5,first=True,align=PP_ALIGN.CENTER)
        x+=cw+gap

import re as _re
ARROW={'comp':(None,'diamond'),'assign':('triangle',None),'real':('triangle',None),'serving':('stealth',None),'flow':('stealth',None),'access':('arrow',None),'assoc':(None,None)}
DASHED={'assoc':'dash','serving':'sysDot','access':'sysDot','flow':'dash'}
def _ln(conn): return conn.line._get_or_add_ln()
def set_ends(conn,head,tail):
    ln=_ln(conn)
    if head: ln.append(ln.makeelement(QN('a:headEnd'),{'type':head,'w':'med','len':'med'}))
    if tail: ln.append(ln.makeelement(QN('a:tailEnd'),{'type':tail,'w':'med','len':'med'}))
def set_dash(conn,val):
    ln=_ln(conn); ln.append(ln.makeelement(QN('a:prstDash'),{'val':val}))
def draw_graph(s,dotpath,bx,by,bw,bh,font=8.0,clusters=None,elabels=False,elsize=7):
    txt=subprocess.run(['dot','-Tplain',dotpath],capture_output=True,text=True).stdout
    nodes={};edges=[];GW=GH=1.0
    for line in txt.splitlines():
        if line.startswith('graph '):
            p=line.split();GW=float(p[2]);GH=float(p[3])
        elif line.startswith('node '):
            m=_re.match(r'node (\S+) (\S+) (\S+) (\S+) (\S+) "(.*?)" \S+ \S+ (\S+) (\S+)',line)
            if m: nodes[m.group(1)]=dict(cx=float(m.group(2)),cy=float(m.group(3)),w=float(m.group(4)),h=float(m.group(5)),lab=m.group(6).replace('\\n','\n'),bd=m.group(7),fl=m.group(8))
        elif line.startswith('edge '):
            p=line.split();n=int(p[3]);rest=p[4+2*n:]
            rt=rest[0] if rest else ''; lx=float(rest[1]) if len(rest)>2 else 0.0; ly=float(rest[2]) if len(rest)>2 else 0.0
            edges.append((p[1],p[2],rt,lx,ly))
    sc=min(bw/GW,bh/GH);offx=bx+(bw-GW*sc)/2;offy=by+(bh-GH*sc)/2
    def place(nd): return offx+(nd['cx']-nd['w']/2)*sc, offy+(GH-nd['cy']-nd['h']/2)*sc, nd['w']*sc, nd['h']*sc
    if clusters:
        for cname,(mem,bg) in clusters.items():
            xs=[];ys=[];xe=[];ye=[]
            for k in mem:
                if k in nodes:
                    x,y,w,h=place(nodes[k]);xs+=[x];ys+=[y];xe+=[x+w];ye+=[y+h]
            if xs:
                pad=0.07
                rect(s,min(xs)-pad,min(ys)-pad-0.17,(max(xe)-min(xs))+2*pad,(max(ye)-min(ys))+2*pad+0.17,RGBColor.from_string(bg),RGBColor(0xBF,0xCA,0xD6),rounded=True,lw=0.75)
                para(tb(s,min(xs)-pad+0.03,min(ys)-pad-0.18,3.5,0.18).text_frame,[(cname,True,RGBColor(0x55,0x60,0x6B))],7,first=True)
    shp={}
    for nid,nd in nodes.items():
        x,y,w,h=place(nd)
        r=rect(s,x,y,w,h,RGBColor.from_string(nd['fl'][1:]),RGBColor.from_string(nd['bd'][1:]),rounded=True,lw=0.75)
        tf=r.text_frame;tf.word_wrap=True;tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        tf.margin_left=Pt(1);tf.margin_right=Pt(1);tf.margin_top=Pt(0);tf.margin_bottom=Pt(0)
        para(tf,[(nd['lab'],False,RGBColor(0x22,0x26,0x2B))],font,first=True,align=PP_ALIGN.CENTER,ls=0.92)
        shp[nid]=r
    for tail,head,rt,lx,ly in edges:
        if tail in shp and head in shp:
            aa=shp[tail];bb=shp[head]
            c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,aa.left+aa.width//2,aa.top+aa.height//2,bb.left+bb.width//2,bb.top+bb.height//2)
            c.line.color.rgb=RGBColor(0x7C,0x86,0x90);c.line.width=Pt(0.75);c.shadow.inherit=False
            he,te=ARROW.get(rt,('arrow',None)); set_ends(c,he,te)
            if rt in DASHED: set_dash(c,DASHED[rt])
            try: c.begin_connect(aa,0);c.end_connect(bb,0)
            except Exception: pass
            if elabels and rt:
                ex=offx+lx*sc; ey=offy+(GH-ly)*sc
                lb=s.shapes.add_textbox(Inches(ex-0.34),Inches(ey-0.085),Inches(0.68),Inches(0.17))
                lb.fill.solid(); lb.fill.fore_color.rgb=RGBColor(0xFF,0xFF,0xFF); lb.line.fill.background(); lb.shadow.inherit=False
                ltf=lb.text_frame; ltf.word_wrap=False
                ltf.margin_left=Pt(0);ltf.margin_right=Pt(0);ltf.margin_top=Pt(0);ltf.margin_bottom=Pt(0)
                para(ltf,[(rt,False,RGBColor(0x6A,0x72,0x7C))],elsize,first=True,align=PP_ALIGN.CENTER)
def setfoot(s):
    for ph in s.placeholders:
        if ph.placeholder_format.type==15:
            ph.text=FOOT
            for r in ph.text_frame.paragraphs[0].runs: r.font.size=Pt(8)
def strip(s,keep=(1,13,15)):  # keep title, slide#, footer
    for ph in list(s.placeholders):
        if ph.placeholder_format.type not in keep:
            ph._element.getparent().remove(ph._element)
def newp(name,title=None,keep=(1,13,15),tsize=20):
    s=prs.slides.add_slide(L[name])
    if title is not None and s.shapes.title is not None:
        s.shapes.title.text=title
        if tsize:
            for r in s.shapes.title.text_frame.paragraphs[0].runs:
                r.font.size=Pt(tsize); r.font.color.rgb=BLACK; r.font.name=F
    strip(s,keep); setfoot(s); return s

OVCL={'1 Tap & Gate':(['d_card','d_wallet','p_nfc','d_nfc','d_pnfc','if_nfc','d_gate','n_si','ss_fw','fn_hs','sv_nfc'],'EAF5E1'),'2 Communication':(['net','p_link'],'FDEFE0'),'3 Backend & Storage':(['n_be','ss_db','fn_store','ar_db','sv_store'],'E3F0F8'),'4 Self-service & Settlement':(['n_as','ar_web'],'F3E9F5')}

# ===== S1 Title (Start) =====
s=prs.slides.add_slide(L['Start'])
if s.shapes.title is not None: s.shapes.title.text="Ticketing & Train Journey Management"
# fill the left object placeholder (idx 10) with info block
for ph in s.placeholders:
    if ph.placeholder_format.idx==10:
        tf=ph.text_frame; tf.clear()
        para(tf,[("Case Study C1 · ArchiTrain · Focus Area 1",False,TUMB)],14,first=True,space=8)
        para(tf,[("Abstract & Detailed ArchiMate Model",False,GREY)],12,space=14)
        para(tf,[("Group C1F1G_",True,BLACK)],12,space=2)
        para(tf,[("Team: Member 1 · Member 2 · Member 3 · …",False,GREY)],11,space=2)
        para(tf,[("TUM School CIT · 2026",False,MUTE)],10,space=2)
setfoot(s)

# ===== S2 Agenda =====
s=newp('Inhalt',"Agenda")
y=1.65
for n,h,d in [("1","Use Case & Scope","ArchiTrain and our focus area F1"),
("2","Abstract Model","The whole enterprise at a glance"),
("3","Detailed Model — Focus Area F1","Business · Application · Technology"),
("4","Reflection","Value · limitations · missing information · challenges")]:
    rect(s,0.6,y,0.5,0.5,TUMB);para(tb(s,0.6,y+0.06,0.5,0.4).text_frame,[(n,True,WHITE)],16,first=True,align=PP_ALIGN.CENTER)
    para(tb(s,1.25,y-0.04,8.2,0.35).text_frame,[(h,True,TUMD)],14,first=True)
    para(tb(s,1.25,y+0.32,8.2,0.35).text_frame,[(d,False,GREY)],11,first=True);y+=0.86

# ===== S3 Use Case =====
s=newp('Inhalt',"ArchiTrain: one architecture, three focus areas")
para(tb(s,0.55,1.62,8.95,0.75).text_frame,[("A modern transit enterprise. Passengers ",False,DARK),("tap on / tap off",True,TUMD),(" with a Go Card, mobile wallet or contactless bank card; the gate opens in under 400 ms — a single touch spanning hardware, networks, backend and payment.",False,DARK)],12,first=True,ls=1.18)
x=0.55
_=0
for tag,h,d,fill,bd,ours in [("F1","Ticketing & Train\nJourney Management","Fare settlement & journey closure",TEC,TEC_B,True),("F2","Train Localization\n& Scheduling","GPS + odometry + track reference points",APPC,APP_B,False),("F3","Malfunction &\nDisruption Management","Threshold detection → human handover",BIZ,BIZ_B,False)]:
    rect(s,x,2.5,2.92,2.0,fill,bd,rounded=True,lw=1.6 if ours else 1.0)
    para(tb(s,x+0.2,2.62,2.55,0.45).text_frame,[(tag,True,TUMD)],20,first=True)
    para(tb(s,x+0.2,3.16,2.55,0.7).text_frame,[(h,True,DARK)],12.5,first=True,ls=1.04)
    para(tb(s,x+0.2,3.95,2.55,0.5).text_frame,[(d,False,GREY)],10,first=True,ls=1.04)
    if ours: para(tb(s,x+0.2,4.22,2.55,0.3).text_frame,[("◀ OUR FOCUS",True,ORANGE)],10,first=True)
    x+=3.05
para(tb(s,0.55,4.66,8.95,0.4).text_frame,[("Material: Sequence Diagram · Interviews (Dr. Marcus Chen · Elena Vasquez · Sarah Mitchell) · Incident Report",False,MUTE)],9.5,first=True)

# ===== S3b Stakeholders & Services =====
s=newp('Inhalt',"Key stakeholders & main services")
# left card: stakeholders
rect(s,0.55,1.62,4.3,3.05,WHITE,LINE,rounded=True);rect(s,0.55,1.62,0.07,3.05,TUMB)
tf=tb(s,0.78,1.74,3.95,3.0).text_frame
para(tf,[("Stakeholders",True,TUMD)],13,first=True,space=7)
for a,b in [("Passenger","Business Actor"),("Customer Support Staff","Business Role"),("Operations Controller","Business Role"),("Train Operator","Business Role"),("Malfunction Management","Business Collaboration")]:
    para(tf,[("• "+a+"  ",False,DARK),("("+b+")",False,MUTE)],11,space=6)
# right card: main services
rect(s,5.15,1.62,4.3,3.05,WHITE,LINE,rounded=True);rect(s,5.15,1.62,0.07,3.05,TEC_B)
tf=tb(s,5.38,1.74,3.95,3.0).text_frame
para(tf,[("Main business services  ",True,TUMD),("(→ serve Passenger)",False,MUTE)],13,first=True,space=7)
for a,b in [("Ticketing & Fare Service","F1"),("Go Card Self-Service","F1"),("On-site Support Service","F1"),("Train Operation & Tracking Service","F2"),("Passenger Information Service","F2 / F3")]:
    para(tf,[("• "+a+"  ",False,DARK),("["+b+"]",False,TUMB)],11,space=6)
para(tb(s,0.55,4.8,8.95,0.4).text_frame,[("Information sources: Elena Vasquez (Head of Fare Systems, F1) · Dr. Marcus Chen (Chief Systems Architect, F2) · Sarah Mitchell (Operations Controller, F3)",False,MUTE)],9,first=True)

# ===== S4 Divider =====
s=newp('Kapiteltrenner',"1   Abstract Model",tsize=None)

# ===== S5 Abstract matrix =====
s=newp('Inhalt',"Abstract model: three domains × three layers")
x0=0.55;lw=1.05;cw=(8.95-lw)/3;y0=1.62;hh=0.34
rect(s,x0,y0,lw,hh,TUMD);para(tb(s,x0,y0+0.06,lw,0.3).text_frame,[("Layer",True,WHITE)],9,first=True,align=PP_ALIGN.CENTER)
for i,c in enumerate(["F1  Ticketing & Journey","F2  Train Localization","F3  Disruption Mgmt"]):
    rect(s,x0+lw+i*cw,y0,cw,hh,TUMD);para(tb(s,x0+lw+i*cw+0.04,y0+0.06,cw-0.08,0.3).text_frame,[(c,True,WHITE)],8.5,first=True,align=PP_ALIGN.CENTER)
rows=[("Business",BIZ,BIZ_B,["Passenger · Customer Support\nHandle Journeys & Fares\n→ Ticketing & Fare Service","Operations Controller · Train Operator\nLocalize & Monitor Trains\n→ Train Operation & Tracking","Malfunction Mgmt (collaboration)\nHandle Delays & Disruptions\n→ Passenger Information Service"]),
("Application",APPC,APP_B,["Access Control Mgmt · JAMS\nFare Management System · Website","LCCU Control Software\nOperational Control Center","Disruption Handling System\n(OCC + PIS + TMS)"]),
("Technology",TEC,TEC_B,["Station Infrastructure · App Server\nJourney & Account Mgmt Backend","LCCU · Central Railway Backend\nRailway Radio","(shared backend)\nPrivate Network"])]
ry=y0+hh;rh=0.95
for name,fill,bd,cells in rows:
    rect(s,x0,ry,lw,rh,bd);para(tb(s,x0,ry+rh/2-0.18,lw,0.4).text_frame,[(name,True,WHITE)],9.5,first=True,align=PP_ALIGN.CENTER)
    for i,cell in enumerate(cells):
        rect(s,x0+lw+i*cw,ry,cw,rh,fill,LINE,lw=0.5);para(tb(s,x0+lw+i*cw+0.08,ry+0.08,cw-0.16,rh-0.12).text_frame,[(cell,False,DARK)],7.8,first=True,ls=1.05)
    ry+=rh
para(tb(s,0.55,ry+0.06,8.95,0.4).text_frame,[("Abstract level: one generalized process per domain · components without internal structure · devices summarized as nodes · communication as networks only.",False,MUTE)],8.5,first=True)

# ===== S6 Divider =====
s=newp('Kapiteltrenner',"2   Detailed Model — Focus Area F1",tsize=None)

# ===== S7 Bridge =====
s=newp('Inhalt',"F1 spans business, application and technology")
y=1.65
for h,fill,bd,d in [("Business",BIZ,BIZ_B,"Conduct Journey · Calculate Fare · Manage Go Card"),("Application",APPC,APP_B,"Access Control · JAMS · Fare Management · Website"),("Technology",TEC,TEC_B,"NFC Reader · Gate · Backend · Private Network · Bank")]:
    rect(s,1.1,y,7.8,0.95,fill,bd,rounded=True,lw=1.2)
    para(tb(s,1.35,y+0.13,7.3,0.4).text_frame,[(h,True,TUMD)],14,first=True);para(tb(s,1.35,y+0.5,7.3,0.4).text_frame,[(d,False,DARK)],11,first=True);y+=1.12
para(tb(s,1.1,y+0.0,7.8,0.4).text_frame,[("Layers are linked by the two backbone relationships: ",False,GREY),("serving",True,ORANGE),(" and ",False,GREY),("realization",True,ORANGE),(".",False,GREY)],11,first=True)

# ===== S8 Business placeholder =====
s=newp('Inhalt',"Business layer")
rect(s,1.5,2.0,7.0,2.5,None,DASH,rounded=True,lw=1.5,dash=True)
para(tb(s,1.5,2.75,7.0,0.6).text_frame,[("Business layer — to be added",True,MUTE)],18,first=True,align=PP_ALIGN.CENTER)
para(tb(s,1.5,3.4,7.0,0.6).text_frame,[("(pending teammate input: actors, processes, services, events)",False,MUTE)],11,first=True,align=PP_ALIGN.CENTER)

# ===== S9 Application placeholder =====
s=newp('Inhalt',"Application layer")
rect(s,1.5,2.0,7.0,2.5,None,DASH,rounded=True,lw=1.5,dash=True)
para(tb(s,1.5,2.75,7.0,0.6).text_frame,[("Application layer — to be added",True,MUTE)],18,first=True,align=PP_ALIGN.CENTER)
para(tb(s,1.5,3.4,7.0,0.6).text_frame,[("(components, application services, data objects)",False,MUTE)],11,first=True,align=PP_ALIGN.CENTER)

# ===== Technology — Overview (bitmap, large) =====
s=newp('Inhalt',"Technology layer — overview & four perspectives"); stepper(s,0)
pic_fit(s,"/Users/ziway/Downloads/tech_overview_cl.png",0.35,1.98,9.3,3.05)
para(tb(s,0.5,5.22,9.0,0.3).text_frame,[("Full technology view grouped into four perspectives — examined one by one next.",False,MUTE)],8.5,first=True,align=PP_ALIGN.CENTER)

# ===== V1 Tap & Gate =====
s=newp('Inhalt',"Perspective 1 — Tap & gate access"); stepper(s,1)
draw_graph(s,"/tmp/tech_v1_tap.dot",0.35,2.0,9.3,2.8,font=8.5,elabels=True)
para(tb(s,0.45,4.88,9.1,0.55).text_frame,[("Passenger taps a Go Card / mobile wallet → NFC Handshake (Path) → reader; NFC Readers realize NFC Proximity Communication Service; Station Gate ⊃ Gate Firmware → NFC Handshake Processing; gate assigned to Access Control Management (< 400 ms).",False,GREY)],8.5,first=True,ls=1.05)

# ===== V2 Communication =====
s=newp('Inhalt',"Perspective 2 — Communication (pattern P01)"); stepper(s,2)
draw_graph(s,"/tmp/tech_v2_comm.dot",0.35,2.0,9.3,2.8,font=8.5,elabels=True)
para(tb(s,0.45,4.88,9.1,0.55).text_frame,[("Card ↔ reader uses NFC Handshake (Path); Gate-Backend Link (Path) is realized by the Private Network (Communication Network) — course pattern P01; the network connects Station Infrastructure, the backend and the Application Server.",False,GREY)],8.5,first=True,ls=1.05)

# ===== V3 Backend & Storage =====
s=newp('Inhalt',"Perspective 3 — Backend storage & data"); stepper(s,3)
draw_graph(s,"/tmp/tech_v3_storage.dot",0.35,2.0,9.3,2.8,font=8.5,elabels=True)
para(tb(s,0.45,4.88,9.1,0.55).text_frame,[("Journey & Account Mgmt Backend ⊃ Database (DBMS) + Journey & Account Database (Artifact); the backend realizes the Data Storage Service serving JAMS & FMS; the artifact realizes Tap Event, Go Card Account, Fare & Zone and Digital Ticket data objects.",False,GREY)],8.5,first=True,ls=1.05)

# ===== V4 Self-service & Settlement =====
s=newp('Inhalt',"Perspective 4 — Self-service & bank settlement"); stepper(s,4)
draw_graph(s,"/tmp/tech_v4_self.dot",0.35,2.0,9.3,2.8,font=8.5,elabels=True)
para(tb(s,0.45,4.88,9.1,0.55).text_frame,[("Application Server ⊃ Website Deployment → realizes the ArchiTrain Website; the app server reads Go Card balances over the Private Network; Bank Payment Service (external) serves the Fare Management System for settlement.",False,GREY)],8.5,first=True,ls=1.05)

# ===== S13 Divider =====
s=newp('Kapiteltrenner',"3   Reflection",tsize=None)

# ===== S14 Reflection placeholder grid =====
s=newp('Inhalt',"Reflection")
quad=[("Business value",TEC_B),("Limitations",ORANGE),("Missing information",TUMB),("Modeling challenges",TUMD)]
gw=4.35;gh=1.35
for i,(h,c) in enumerate(quad):
    cx=0.55+(i%2)*(gw+0.25); cy=1.62+(i//2)*(gh+0.18)
    rect(s,cx,cy,gw,gh,WHITE,LINE,rounded=True);rect(s,cx,cy,0.07,gh,c)
    para(tb(s,cx+0.2,cy+0.14,gw-0.35,0.4).text_frame,[(h,True,TUMD)],12.5,first=True)
    para(tb(s,cx+0.2,cy+0.6,gw-0.35,0.6).text_frame,[("[ to be completed with the team ]",False,MUTE)],10,first=True)

# ===== S15 Conclusion =====
s=newp('Inhalt',"Conclusion")
tf=tb(s,0.6,1.65,8.8,2.2).text_frame
para(tf,[("• The abstract model covers the whole enterprise (F1 + F2 + F3) with structure and behaviour",False,DARK)],12,first=True,space=10)
para(tf,[("• The F1 detailed model is developed across business, application and technology — communication explicitly modelled",False,DARK)],12,space=10)
para(tf,[("• The vertical chain Node → Component → App Service → Business Process → Business Service → Passenger is derivable",False,DARK)],12,space=10)
rect(s,0.6,4.2,8.8,0.8,TUMB)
para(tb(s,0.9,4.36,8.2,0.5).text_frame,[("Thank you — questions & discussion welcome.",True,WHITE)],15,first=True)

prs.save("/Users/ziway/Downloads/C1F1G_ArchiTrain_Presentation.pptx")
print("saved", len(prs.slides._sldIdLst),"slides on TUM template")
